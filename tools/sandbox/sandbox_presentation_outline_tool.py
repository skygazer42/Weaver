"""
Sandbox Presentation Outline Tool for LLM-based PPT outline generation.

This module uses LLM to generate structured presentation outlines:
- Analyze topic and generate logical slide structure
- Create title, content, and speaker notes for each slide
- Support multiple presentation styles (business, educational, technical)
- Export outline to JSON or directly create presentation

Similar to Manus's sb_presentation_outline_tool.py but adapted for Weaver.

Usage:
    from tools.sandbox.sandbox_presentation_outline_tool import build_presentation_outline_tools

    tools = build_presentation_outline_tools(thread_id="thread_123")
"""

from __future__ import annotations

import asyncio
import json
import logging
import time
from typing import Any, Dict, List, Optional, Literal

from langchain_core.tools import BaseTool
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_openai import ChatOpenAI
from pydantic import BaseModel, Field

from common.config import settings

logger = logging.getLogger(__name__)


# Presentation styles
PresentationStyle = Literal["business", "educational", "technical", "creative", "minimal"]


class SlideOutline(BaseModel):
    """Structure for a single slide outline."""
    slide_number: int = Field(description="Slide number (1-based)")
    title: str = Field(description="Slide title")
    layout: str = Field(description="Slide layout type")
    content: List[str] = Field(description="Bullet points or content items")
    speaker_notes: str = Field(default="", description="Speaker notes for this slide")
    has_image: bool = Field(default=False, description="Whether this slide should have an image")
    image_suggestion: str = Field(default="", description="Suggested image description if has_image is True")


class PresentationOutline(BaseModel):
    """Complete presentation outline structure."""
    title: str = Field(description="Presentation title")
    subtitle: str = Field(default="", description="Presentation subtitle")
    author: str = Field(default="", description="Author name")
    total_slides: int = Field(description="Total number of slides")
    estimated_duration_minutes: int = Field(description="Estimated presentation duration")
    target_audience: str = Field(description="Target audience description")
    key_takeaways: List[str] = Field(description="Key takeaways from the presentation")
    slides: List[SlideOutline] = Field(description="List of slide outlines")


OUTLINE_SYSTEM_PROMPT = """You are an expert presentation designer. Your task is to create a structured outline for a presentation.

## Guidelines:

1. **Structure**: Create a logical flow with clear sections:
   - Title slide
   - Agenda/Overview (optional)
   - Main content sections
   - Summary/Key takeaways
   - Q&A or Call to action

2. **Content per slide**:
   - 3-5 bullet points maximum
   - Concise, impactful phrases
   - Action-oriented language

3. **Layouts**: Use appropriate layouts:
   - "title" - Title slide with subtitle
   - "title_content" - Title with bullet points
   - "section" - Section header
   - "two_content" - Side-by-side comparison
   - "blank" - For images or custom content

4. **Speaker Notes**: Add helpful notes for the presenter

5. **Images**: Suggest images for visual slides

## Style Guidelines:

- **business**: Professional, data-driven, executive summary style
- **educational**: Teaching-focused, step-by-step, with examples
- **technical**: Detailed, code/diagram heavy, precise terminology
- **creative**: Visually engaging, storytelling, minimal text
- **minimal**: Clean, few words, impact-focused

Respond with a valid JSON object matching the PresentationOutline schema.
"""


def _get_event_emitter(thread_id: str):
    """Get event emitter for a thread."""
    from agent.core.events import get_emitter_sync

    return get_emitter_sync(thread_id)


class _PresentationOutlineBaseTool(BaseTool):
    """Base class for presentation outline tools."""

    thread_id: str = "default"
    emit_events: bool = True

    def _emit_event(self, event_type: str, data: Dict[str, Any]) -> None:
        """Emit an event."""
        if not self.emit_events:
            return
        emitter = _get_event_emitter(self.thread_id)
        if emitter:
            try:
                emitter.emit_sync(event_type, data)
            except Exception as e:
                logger.warning(f"[presentation_outline] Failed to emit event: {e}")

    def _emit_tool_start(self, action: str, args: Dict[str, Any]) -> float:
        """Emit tool start event."""
        start_time = time.time()
        self._emit_event("tool_start", {
            "tool": self.name,
            "action": action,
            "args": args,
            "thread_id": self.thread_id,
        })
        return start_time

    def _emit_tool_result(
        self,
        action: str,
        result: Dict[str, Any],
        start_time: float,
        success: bool = True,
    ) -> None:
        """Emit tool result event."""
        duration_ms = (time.time() - start_time) * 1000
        self._emit_event("tool_result", {
            "tool": self.name,
            "action": action,
            "success": success,
            "duration_ms": round(duration_ms, 2),
        })

    def _get_llm(self) -> ChatOpenAI:
        """Get LLM for outline generation."""
        model = settings.reasoning_model or "gpt-4o-mini"
        params = {
            "model": model,
            "temperature": 0.7,
            "api_key": settings.openai_api_key,
            "timeout": settings.openai_timeout or 60,
        }
        if settings.openai_base_url:
            params["base_url"] = settings.openai_base_url
        return ChatOpenAI(**params)


class GenerateOutlineInput(BaseModel):
    """Input for generate_outline."""
    topic: str = Field(description="The presentation topic or title")
    num_slides: int = Field(
        default=10,
        ge=3,
        le=30,
        description="Approximate number of slides (3-30)"
    )
    style: PresentationStyle = Field(
        default="business",
        description="Presentation style: business, educational, technical, creative, minimal"
    )
    target_audience: str = Field(
        default="general",
        description="Target audience description"
    )
    duration_minutes: int = Field(
        default=15,
        description="Target presentation duration in minutes"
    )
    include_images: bool = Field(
        default=True,
        description="Whether to suggest images for slides"
    )
    additional_context: str = Field(
        default="",
        description="Additional context or requirements"
    )


class GenerateOutlineTool(_PresentationOutlineBaseTool):
    """Generate a presentation outline using LLM."""

    name: str = "generate_presentation_outline"
    description: str = (
        "Generate a structured presentation outline using AI. "
        "Provide the topic, number of slides, and style preference. "
        "Returns a complete outline with titles, content, and speaker notes."
    )
    args_schema: type[BaseModel] = GenerateOutlineInput

    def _run(
        self,
        topic: str,
        num_slides: int = 10,
        style: PresentationStyle = "business",
        target_audience: str = "general",
        duration_minutes: int = 15,
        include_images: bool = True,
        additional_context: str = "",
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("generate_outline", {
            "topic": topic,
            "num_slides": num_slides,
            "style": style,
        })

        try:
            llm = self._get_llm()

            # Build the prompt
            user_prompt = f"""Create a presentation outline for:

**Topic**: {topic}
**Number of slides**: approximately {num_slides}
**Style**: {style}
**Target audience**: {target_audience}
**Duration**: {duration_minutes} minutes
**Include image suggestions**: {include_images}

{f"**Additional context**: {additional_context}" if additional_context else ""}

Generate a complete JSON outline following the PresentationOutline schema.
"""

            messages = [
                SystemMessage(content=OUTLINE_SYSTEM_PROMPT),
                HumanMessage(content=user_prompt),
            ]

            # Get structured output
            response = llm.with_structured_output(PresentationOutline).invoke(messages)

            outline_dict = response.model_dump()

            result = {
                "success": True,
                "message": f"Generated outline for '{topic}' with {len(response.slides)} slides",
                "outline": outline_dict,
            }

            self._emit_tool_result("generate_outline", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("generate_outline", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class OutlineToSlidesInput(BaseModel):
    """Input for outline_to_slides."""
    outline: Dict[str, Any] = Field(description="The presentation outline (from generate_outline)")
    file_path: str = Field(description="Output path for the PPTX file")


class OutlineToSlidesTool(_PresentationOutlineBaseTool):
    """Convert an outline to an actual presentation file."""

    name: str = "outline_to_presentation"
    description: str = (
        "Convert a presentation outline (from generate_outline) to an actual PPTX file. "
        "Creates slides with the specified content and layout."
    )
    args_schema: type[BaseModel] = OutlineToSlidesInput

    def _get_sandbox(self):
        """Get the E2B sandbox instance."""
        try:
            from tools.sandbox.sandbox_browser_session import sandbox_browser_sessions
            session = sandbox_browser_sessions.get(self.thread_id)
            if session and hasattr(session, "_handles") and session._handles:
                return session._handles.sandbox
        except Exception:
            pass
        return None

    def _run(
        self,
        outline: Dict[str, Any],
        file_path: str,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("outline_to_slides", {
            "file_path": file_path,
            "slide_count": len(outline.get("slides", [])),
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized. Start sandbox browser first.")

            # Ensure python-pptx is installed
            sandbox.commands.run("pip install python-pptx", timeout=60)

            workspace_path = "/workspace"
            full_path = f"{workspace_path}/{file_path.lstrip('/')}"

            # Create parent directories
            parent_dir = "/".join(full_path.split("/")[:-1])
            if parent_dir:
                sandbox.commands.run(f"mkdir -p {parent_dir}")

            # Serialize outline to JSON for Python script
            outline_json = json.dumps(outline, ensure_ascii=False)

            python_code = f'''
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

outline = json.loads("""{outline_json}""")

prs = Presentation()

# Layout mapping
layout_map = {{
    "title": 0,
    "title_content": 1,
    "section": 2,
    "two_content": 3,
    "blank": 6,
}}

# Create title slide
title_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(title_layout)
title_slide.shapes.title.text = outline.get("title", "Presentation")
if len(title_slide.placeholders) > 1:
    subtitle = outline.get("subtitle", "")
    if outline.get("author"):
        subtitle += f"\\n{{outline['author']}}"
    title_slide.placeholders[1].text = subtitle

# Create content slides
for slide_data in outline.get("slides", []):
    layout_name = slide_data.get("layout", "title_content")
    layout_idx = layout_map.get(layout_name, 1)

    try:
        slide_layout = prs.slide_layouts[layout_idx]
    except:
        slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "")

    # Set content
    content = slide_data.get("content", [])
    if content and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for i, item in enumerate(content):
            if i == 0:
                tf.paragraphs[0].text = str(item)
            else:
                p = tf.add_paragraph()
                p.text = str(item)
                p.level = 0

    # Add speaker notes
    notes = slide_data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

prs.save("{full_path}")
print(f"SUCCESS: Created {{len(prs.slides)}} slides")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=60)

            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to create presentation: {result.stderr}")

            # Extract slide count
            slide_count = len(outline.get("slides", [])) + 1  # +1 for title slide

            result = {
                "success": True,
                "message": f"Created presentation with {slide_count} slides",
                "path": file_path,
                "slide_count": slide_count,
            }

            self._emit_tool_result("outline_to_slides", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("outline_to_slides", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class RefineOutlineInput(BaseModel):
    """Input for refine_outline."""
    outline: Dict[str, Any] = Field(description="The current presentation outline")
    feedback: str = Field(description="Feedback or changes to apply")


class RefineOutlineTool(_PresentationOutlineBaseTool):
    """Refine an existing outline based on feedback."""

    name: str = "refine_presentation_outline"
    description: str = (
        "Refine an existing presentation outline based on feedback. "
        "Provide the current outline and the changes you want to make."
    )
    args_schema: type[BaseModel] = RefineOutlineInput

    def _run(
        self,
        outline: Dict[str, Any],
        feedback: str,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("refine_outline", {
            "feedback": feedback[:100] + "..." if len(feedback) > 100 else feedback,
        })

        try:
            llm = self._get_llm()

            user_prompt = f"""Here is the current presentation outline:

```json
{json.dumps(outline, indent=2, ensure_ascii=False)}
```

**Requested changes**:
{feedback}

Please modify the outline according to the feedback and return the complete updated outline as JSON.
"""

            messages = [
                SystemMessage(content=OUTLINE_SYSTEM_PROMPT),
                HumanMessage(content=user_prompt),
            ]

            response = llm.with_structured_output(PresentationOutline).invoke(messages)
            outline_dict = response.model_dump()

            result = {
                "success": True,
                "message": "Outline refined successfully",
                "outline": outline_dict,
            }

            self._emit_tool_result("refine_outline", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("refine_outline", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class ExpandSlideInput(BaseModel):
    """Input for expand_slide."""
    topic: str = Field(description="The slide topic to expand")
    context: str = Field(default="", description="Context from surrounding slides")
    style: PresentationStyle = Field(default="business", description="Presentation style")


class ExpandSlideTool(_PresentationOutlineBaseTool):
    """Generate detailed content for a single slide."""

    name: str = "expand_slide_content"
    description: str = (
        "Generate detailed content for a single slide topic. "
        "Useful for expanding brief outline points into full slide content."
    )
    args_schema: type[BaseModel] = ExpandSlideInput

    def _run(
        self,
        topic: str,
        context: str = "",
        style: PresentationStyle = "business",
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("expand_slide", {"topic": topic})

        try:
            llm = self._get_llm()

            user_prompt = f"""Generate detailed slide content for:

**Topic**: {topic}
**Style**: {style}
{f"**Context**: {context}" if context else ""}

Provide:
1. A compelling title
2. 4-5 bullet points with supporting details
3. Speaker notes (2-3 sentences)
4. An image suggestion if appropriate

Return as a SlideOutline JSON object.
"""

            messages = [
                SystemMessage(content="You are an expert at creating presentation slide content. Generate clear, impactful content."),
                HumanMessage(content=user_prompt),
            ]

            response = llm.with_structured_output(SlideOutline).invoke(messages)
            slide_dict = response.model_dump()

            result = {
                "success": True,
                "slide": slide_dict,
            }

            self._emit_tool_result("expand_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("expand_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


def build_presentation_outline_tools(
    thread_id: str,
    emit_events: bool = True,
) -> List[BaseTool]:
    """
    Build presentation outline tools for a thread.

    Args:
        thread_id: Thread/conversation ID
        emit_events: Whether to emit events

    Returns:
        List of presentation outline tools
    """
    return [
        GenerateOutlineTool(thread_id=thread_id, emit_events=emit_events),
        OutlineToSlidesTool(thread_id=thread_id, emit_events=emit_events),
        RefineOutlineTool(thread_id=thread_id, emit_events=emit_events),
        ExpandSlideTool(thread_id=thread_id, emit_events=emit_events),
    ]
