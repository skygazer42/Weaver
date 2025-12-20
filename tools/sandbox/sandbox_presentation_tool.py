"""
Sandbox Presentation Tool for E2B Sandbox PowerPoint Operations.

This module provides presentation generation capabilities in an E2B sandbox:
- Create PowerPoint (.pptx) presentations
- Add slides with various layouts
- Add text, images, shapes, and tables
- Apply themes and formatting
- Export to PDF

Similar to Manus's sb_presentation_tool.py but adapted for Weaver's E2B integration.

Usage:
    from tools.sandbox.sandbox_presentation_tool import build_sandbox_presentation_tools

    tools = build_sandbox_presentation_tools(thread_id="thread_123")
"""

from __future__ import annotations

import asyncio
import base64
import json
import logging
import time
from typing import Any, Dict, List, Optional, Union

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


# Slide layout types
SLIDE_LAYOUTS = {
    "title": 0,           # Title slide
    "title_content": 1,   # Title and Content
    "section": 2,         # Section Header
    "two_content": 3,     # Two Content
    "comparison": 4,      # Comparison
    "title_only": 5,      # Title Only
    "blank": 6,           # Blank
    "content_caption": 7, # Content with Caption
    "picture_caption": 8, # Picture with Caption
}


def _get_sandbox_session(thread_id: str):
    """Get sandbox session for a thread."""
    try:
        from tools.sandbox.sandbox_browser_session import sandbox_browser_sessions
        return sandbox_browser_sessions.get(thread_id)
    except ImportError:
        return None


def _get_event_emitter(thread_id: str):
    """Get event emitter for a thread."""
    try:
        from agent.events import get_emitter_sync
        return get_emitter_sync(thread_id)
    except ImportError:
        return None


class _SandboxPresentationBaseTool(BaseTool):
    """Base class for sandbox presentation tools."""

    thread_id: str = "default"
    emit_events: bool = True
    workspace_path: str = "/workspace"

    def _get_sandbox(self):
        """Get the E2B sandbox instance."""
        session = _get_sandbox_session(self.thread_id)
        if session and hasattr(session, "_handles") and session._handles:
            return session._handles.sandbox
        return None

    def _emit_event(self, event_type: str, data: Dict[str, Any]) -> None:
        """Emit an event."""
        if not self.emit_events:
            return
        emitter = _get_event_emitter(self.thread_id)
        if emitter:
            try:
                loop = asyncio.new_event_loop()
                try:
                    loop.run_until_complete(emitter.emit(event_type, data))
                finally:
                    loop.close()
            except Exception as e:
                logger.warning(f"[sandbox_presentation] Failed to emit event: {e}")

    def _emit_tool_start(self, action: str, args: Dict[str, Any]) -> float:
        """Emit tool start event."""
        start_time = time.time()
        self._emit_event("tool_start", {
            "tool": self.name,
            "action": action,
            "args": args,
            "thread_id": self.thread_id,
        })
        return start_time

    def _emit_tool_result(
        self,
        action: str,
        result: Dict[str, Any],
        start_time: float,
        success: bool = True,
    ) -> None:
        """Emit tool result event."""
        duration_ms = (time.time() - start_time) * 1000
        self._emit_event("tool_result", {
            "tool": self.name,
            "action": action,
            "success": success,
            "duration_ms": round(duration_ms, 2),
        })

    def _ensure_pptx(self, sandbox) -> bool:
        """Ensure python-pptx is installed in sandbox."""
        try:
            result = sandbox.commands.run("pip show python-pptx", timeout=30)
            if result.exit_code != 0:
                logger.info("[sandbox_presentation] Installing python-pptx...")
                install_result = sandbox.commands.run(
                    "pip install python-pptx Pillow",
                    timeout=120
                )
                return install_result.exit_code == 0
            return True
        except Exception as e:
            logger.warning(f"[sandbox_presentation] Failed to check/install python-pptx: {e}")
            return False


class CreatePresentationInput(BaseModel):
    """Input for create_presentation."""
    file_path: str = Field(
        description="Path for the presentation file (e.g., 'presentations/demo.pptx')"
    )
    title: str = Field(
        default="",
        description="Title for the first slide"
    )
    subtitle: str = Field(
        default="",
        description="Subtitle for the first slide"
    )


class SandboxCreatePresentationTool(_SandboxPresentationBaseTool):
    """Create a new PowerPoint presentation."""

    name: str = "sandbox_create_presentation"
    description: str = (
        "Create a new PowerPoint (.pptx) presentation file. "
        "Optionally set the title slide content. "
        "Path must be relative to /workspace."
    )
    args_schema: type[BaseModel] = CreatePresentationInput

    def _run(
        self,
        file_path: str,
        title: str = "",
        subtitle: str = "",
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("create_presentation", {"file_path": file_path})

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized. Start sandbox browser first.")

            # Ensure python-pptx is installed
            if not self._ensure_pptx(sandbox):
                return {"success": False, "error": "Failed to install python-pptx"}

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            # Create parent directories
            parent_dir = "/".join(full_path.split("/")[:-1])
            if parent_dir:
                sandbox.commands.run(f"mkdir -p {parent_dir}")

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Add title slide if title provided
title_text = """{title}"""
subtitle_text = """{subtitle}"""

if title_text:
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title_text
    if subtitle_text:
        subtitle_shape.text = subtitle_text

prs.save("{full_path}")
print("SUCCESS")
print(f"Slides: {{len(prs.slides)}}")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to create presentation: {result.stderr}")

            result = {
                "success": True,
                "message": f"Presentation created: {file_path}",
                "path": file_path,
                "has_title_slide": bool(title),
            }

            self._emit_tool_result("create_presentation", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("create_presentation", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class AddSlideInput(BaseModel):
    """Input for add_slide."""
    file_path: str = Field(description="Path to the presentation file")
    layout: str = Field(
        default="title_content",
        description="Slide layout: 'title', 'title_content', 'section', 'two_content', 'comparison', 'title_only', 'blank'"
    )
    title: str = Field(default="", description="Slide title")
    content: Optional[str] = Field(default=None, description="Main content text (supports bullet points with newlines)")
    notes: Optional[str] = Field(default=None, description="Speaker notes")


class SandboxAddSlideTool(_SandboxPresentationBaseTool):
    """Add a slide to a presentation."""

    name: str = "sandbox_add_slide"
    description: str = (
        "Add a new slide to an existing presentation. "
        "Choose from various layouts: 'title', 'title_content', 'section', 'two_content', 'blank'. "
        "Content can include bullet points (separated by newlines)."
    )
    args_schema: type[BaseModel] = AddSlideInput

    def _run(
        self,
        file_path: str,
        layout: str = "title_content",
        title: str = "",
        content: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("add_slide", {
            "file_path": file_path,
            "layout": layout,
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            # Map layout name to index
            layout_idx = SLIDE_LAYOUTS.get(layout, 1)

            # Escape content for Python
            content_escaped = (content or "").replace('"""', '\\"\\"\\"').replace("\\", "\\\\")
            notes_escaped = (notes or "").replace('"""', '\\"\\"\\"').replace("\\", "\\\\")

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation("{full_path}")

layout_idx = {layout_idx}
slide_layout = prs.slide_layouts[layout_idx]
slide = prs.slides.add_slide(slide_layout)

# Set title
title_text = """{title}"""
if slide.shapes.title and title_text:
    slide.shapes.title.text = title_text

# Set content
content_text = """{content_escaped}"""
if content_text and len(slide.placeholders) > 1:
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    lines = content_text.split("\\n")
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line.strip()
        else:
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 0

# Add notes
notes_text = """{notes_escaped}"""
if notes_text:
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

prs.save("{full_path}")
print("SUCCESS")
print(f"Total slides: {{len(prs.slides)}}")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to add slide: {result.stderr}")

            # Extract slide count
            slide_count = 0
            for line in result.stdout.split("\n"):
                if "Total slides:" in line:
                    try:
                        slide_count = int(line.split(":")[-1].strip())
                    except Exception:
                        pass

            result = {
                "success": True,
                "message": f"Slide added (layout: {layout})",
                "path": file_path,
                "slide_number": slide_count,
                "layout": layout,
            }

            self._emit_tool_result("add_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("add_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class AddImageToSlideInput(BaseModel):
    """Input for add_image_to_slide."""
    file_path: str = Field(description="Path to the presentation file")
    slide_number: int = Field(description="Slide number (1-based)")
    image_path: str = Field(description="Path to the image file in sandbox")
    left: float = Field(default=1.0, description="Left position in inches")
    top: float = Field(default=2.0, description="Top position in inches")
    width: Optional[float] = Field(default=None, description="Width in inches (auto if not set)")
    height: Optional[float] = Field(default=None, description="Height in inches (auto if not set)")


class SandboxAddImageToSlideTool(_SandboxPresentationBaseTool):
    """Add an image to a slide."""

    name: str = "sandbox_add_image_to_slide"
    description: str = (
        "Add an image to a specific slide. "
        "Position and size are specified in inches. "
        "Image must exist in the sandbox filesystem."
    )
    args_schema: type[BaseModel] = AddImageToSlideInput

    def _run(
        self,
        file_path: str,
        slide_number: int,
        image_path: str,
        left: float = 1.0,
        top: float = 2.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("add_image_to_slide", {
            "file_path": file_path,
            "slide_number": slide_number,
            "image_path": image_path,
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"
            full_image_path = f"{self.workspace_path}/{image_path.lstrip('/')}"

            width_arg = f"Inches({width})" if width else "None"
            height_arg = f"Inches({height})" if height else "None"

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches

prs = Presentation("{full_path}")

slide_idx = {slide_number} - 1
if slide_idx < 0 or slide_idx >= len(prs.slides):
    print("ERROR: Invalid slide number")
    exit(1)

slide = prs.slides[slide_idx]

left = Inches({left})
top = Inches({top})
width = {width_arg}
height = {height_arg}

slide.shapes.add_picture("{full_image_path}", left, top, width, height)

prs.save("{full_path}")
print("SUCCESS")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to add image: {result.stderr}")

            result = {
                "success": True,
                "message": f"Image added to slide {slide_number}",
                "path": file_path,
                "slide_number": slide_number,
                "image_path": image_path,
            }

            self._emit_tool_result("add_image_to_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("add_image_to_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class AddTableToSlideInput(BaseModel):
    """Input for add_table_to_slide."""
    file_path: str = Field(description="Path to the presentation file")
    slide_number: int = Field(description="Slide number (1-based)")
    data: List[List[str]] = Field(description="2D array of table data (first row as headers)")
    left: float = Field(default=1.0, description="Left position in inches")
    top: float = Field(default=2.0, description="Top position in inches")
    width: float = Field(default=8.0, description="Table width in inches")
    height: float = Field(default=3.0, description="Table height in inches")


class SandboxAddTableToSlideTool(_SandboxPresentationBaseTool):
    """Add a table to a slide."""

    name: str = "sandbox_add_table_to_slide"
    description: str = (
        "Add a table to a specific slide. "
        "Provide data as a 2D array where the first row is headers."
    )
    args_schema: type[BaseModel] = AddTableToSlideInput

    def _run(
        self,
        file_path: str,
        slide_number: int,
        data: List[List[str]],
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("add_table_to_slide", {
            "file_path": file_path,
            "slide_number": slide_number,
            "rows": len(data),
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"
            data_json = json.dumps(data)

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
import json

prs = Presentation("{full_path}")

slide_idx = {slide_number} - 1
if slide_idx < 0 or slide_idx >= len(prs.slides):
    print("ERROR: Invalid slide number")
    exit(1)

slide = prs.slides[slide_idx]
data = json.loads('{data_json}')

rows = len(data)
cols = len(data[0]) if data else 0

table = slide.shapes.add_table(
    rows, cols,
    Inches({left}), Inches({top}),
    Inches({width}), Inches({height})
).table

# Populate table
for row_idx, row_data in enumerate(data):
    for col_idx, cell_value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = str(cell_value)

        # Make header row bold
        if row_idx == 0:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

prs.save("{full_path}")
print("SUCCESS")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to add table: {result.stderr}")

            result = {
                "success": True,
                "message": f"Table added to slide {slide_number}",
                "path": file_path,
                "slide_number": slide_number,
                "rows": len(data),
                "columns": len(data[0]) if data else 0,
            }

            self._emit_tool_result("add_table_to_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("add_table_to_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class AddShapeToSlideInput(BaseModel):
    """Input for add_shape_to_slide."""
    file_path: str = Field(description="Path to the presentation file")
    slide_number: int = Field(description="Slide number (1-based)")
    shape_type: str = Field(
        description="Shape type: 'rectangle', 'oval', 'rounded_rectangle', 'arrow_right', 'arrow_left'"
    )
    left: float = Field(description="Left position in inches")
    top: float = Field(description="Top position in inches")
    width: float = Field(description="Width in inches")
    height: float = Field(description="Height in inches")
    text: Optional[str] = Field(default=None, description="Text inside the shape")
    fill_color: Optional[str] = Field(default=None, description="Fill color (hex, e.g., '4472C4')")


class SandboxAddShapeToSlideTool(_SandboxPresentationBaseTool):
    """Add a shape to a slide."""

    name: str = "sandbox_add_shape_to_slide"
    description: str = (
        "Add a shape (rectangle, oval, arrow, etc.) to a slide. "
        "Optionally add text inside and set fill color."
    )
    args_schema: type[BaseModel] = AddShapeToSlideInput

    def _run(
        self,
        file_path: str,
        slide_number: int,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("add_shape_to_slide", {
            "file_path": file_path,
            "slide_number": slide_number,
            "shape_type": shape_type,
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            shape_map = {
                "rectangle": "RECTANGLE",
                "oval": "OVAL",
                "rounded_rectangle": "ROUNDED_RECTANGLE",
                "arrow_right": "RIGHT_ARROW",
                "arrow_left": "LEFT_ARROW",
                "triangle": "ISOSCELES_TRIANGLE",
                "diamond": "DIAMOND",
                "pentagon": "PENTAGON",
                "hexagon": "HEXAGON",
                "star": "STAR_5_POINT",
            }
            shape_const = shape_map.get(shape_type.lower(), "RECTANGLE")

            text_escaped = (text or "").replace('"', '\\"')
            color_code = f'RgbColor.from_string("{fill_color}")' if fill_color else "None"

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RgbColor

prs = Presentation("{full_path}")

slide_idx = {slide_number} - 1
if slide_idx < 0 or slide_idx >= len(prs.slides):
    print("ERROR: Invalid slide number")
    exit(1)

slide = prs.slides[slide_idx]

shape = slide.shapes.add_shape(
    MSO_SHAPE.{shape_const},
    Inches({left}), Inches({top}),
    Inches({width}), Inches({height})
)

# Set fill color
fill_color = {color_code}
if fill_color:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

# Add text
text = "{text_escaped}"
if text:
    shape.text = text

prs.save("{full_path}")
print("SUCCESS")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to add shape: {result.stderr}")

            result = {
                "success": True,
                "message": f"Shape '{shape_type}' added to slide {slide_number}",
                "path": file_path,
                "slide_number": slide_number,
                "shape_type": shape_type,
            }

            self._emit_tool_result("add_shape_to_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("add_shape_to_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class GetPresentationInfoInput(BaseModel):
    """Input for get_presentation_info."""
    file_path: str = Field(description="Path to the presentation file")


class SandboxGetPresentationInfoTool(_SandboxPresentationBaseTool):
    """Get information about a presentation."""

    name: str = "sandbox_get_presentation_info"
    description: str = (
        "Get information about a presentation including slide count, "
        "slide titles, and dimensions."
    )
    args_schema: type[BaseModel] = GetPresentationInfoInput

    def _run(self, file_path: str) -> Dict[str, Any]:
        start_time = self._emit_tool_start("get_presentation_info", {"file_path": file_path})

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            python_code = f'''
from pptx import Presentation
from pptx.util import Inches
import json

prs = Presentation("{full_path}")

slides_info = []
for idx, slide in enumerate(prs.slides):
    slide_info = {{"number": idx + 1, "title": ""}}
    if slide.shapes.title:
        slide_info["title"] = slide.shapes.title.text
    slide_info["shape_count"] = len(slide.shapes)
    slides_info.append(slide_info)

info = {{
    "slide_count": len(prs.slides),
    "width_inches": prs.slide_width.inches,
    "height_inches": prs.slide_height.inches,
    "slides": slides_info
}}

print(json.dumps(info))
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)

            try:
                info = json.loads(result.stdout.strip())
            except Exception:
                raise RuntimeError(f"Failed to parse presentation info: {result.stderr}")

            result = {
                "success": True,
                "path": file_path,
                **info
            }

            self._emit_tool_result("get_presentation_info", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("get_presentation_info", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class UpdateSlideInput(BaseModel):
    """Input for update_slide."""
    file_path: str = Field(description="Path to the presentation file")
    slide_number: int = Field(description="Slide number (1-based)")
    title: Optional[str] = Field(default=None, description="New title (None to keep existing)")
    content: Optional[str] = Field(default=None, description="New content (None to keep existing)")


class SandboxUpdateSlideTool(_SandboxPresentationBaseTool):
    """Update an existing slide's content."""

    name: str = "sandbox_update_slide"
    description: str = (
        "Update the title and/or content of an existing slide. "
        "Only provided fields will be updated."
    )
    args_schema: type[BaseModel] = UpdateSlideInput

    def _run(
        self,
        file_path: str,
        slide_number: int,
        title: Optional[str] = None,
        content: Optional[str] = None,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("update_slide", {
            "file_path": file_path,
            "slide_number": slide_number,
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            title_arg = f'"""{title}"""' if title is not None else "None"
            content_escaped = (content or "").replace('"""', '\\"\\"\\"')
            content_arg = f'"""{content_escaped}"""' if content is not None else "None"

            python_code = f'''
from pptx import Presentation

prs = Presentation("{full_path}")

slide_idx = {slide_number} - 1
if slide_idx < 0 or slide_idx >= len(prs.slides):
    print("ERROR: Invalid slide number")
    exit(1)

slide = prs.slides[slide_idx]

new_title = {title_arg}
new_content = {content_arg}

if new_title is not None and slide.shapes.title:
    slide.shapes.title.text = new_title

if new_content is not None and len(slide.placeholders) > 1:
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    lines = new_content.split("\\n")
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line.strip()
        else:
            p = tf.add_paragraph()
            p.text = line.strip()

prs.save("{full_path}")
print("SUCCESS")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to update slide: {result.stderr}")

            result = {
                "success": True,
                "message": f"Slide {slide_number} updated",
                "path": file_path,
                "slide_number": slide_number,
            }

            self._emit_tool_result("update_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("update_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


class DeleteSlideInput(BaseModel):
    """Input for delete_slide."""
    file_path: str = Field(description="Path to the presentation file")
    slide_number: int = Field(description="Slide number to delete (1-based)")


class SandboxDeleteSlideTool(_SandboxPresentationBaseTool):
    """Delete a slide from a presentation."""

    name: str = "sandbox_delete_slide"
    description: str = "Delete a slide from a presentation by slide number (1-based)."
    args_schema: type[BaseModel] = DeleteSlideInput

    def _run(
        self,
        file_path: str,
        slide_number: int,
    ) -> Dict[str, Any]:
        start_time = self._emit_tool_start("delete_slide", {
            "file_path": file_path,
            "slide_number": slide_number,
        })

        try:
            sandbox = self._get_sandbox()
            if not sandbox:
                raise RuntimeError("Sandbox not initialized.")

            full_path = f"{self.workspace_path}/{file_path.lstrip('/')}"

            python_code = f'''
from pptx import Presentation

prs = Presentation("{full_path}")

slide_idx = {slide_number} - 1
if slide_idx < 0 or slide_idx >= len(prs.slides):
    print("ERROR: Invalid slide number")
    exit(1)

# Get the slide to delete
slide_id = prs.slides._sldIdLst[slide_idx].rId
prs.part.drop_rel(slide_id)
del prs.slides._sldIdLst[slide_idx]

prs.save("{full_path}")
print("SUCCESS")
print(f"Remaining slides: {{len(prs.slides)}}")
'''
            result = sandbox.commands.run(f"python3 -c '{python_code}'", timeout=30)
            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"Failed to delete slide: {result.stderr}")

            result = {
                "success": True,
                "message": f"Slide {slide_number} deleted",
                "path": file_path,
            }

            self._emit_tool_result("delete_slide", result, start_time, True)
            return result

        except Exception as e:
            self._emit_tool_result("delete_slide", {"error": str(e)}, start_time, False)
            return {"success": False, "error": str(e)}


def build_sandbox_presentation_tools(
    thread_id: str,
    emit_events: bool = True,
) -> List[BaseTool]:
    """
    Build sandbox presentation tools for a thread.

    Args:
        thread_id: Thread/conversation ID
        emit_events: Whether to emit events

    Returns:
        List of presentation tools
    """
    return [
        SandboxCreatePresentationTool(thread_id=thread_id, emit_events=emit_events),
        SandboxAddSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxUpdateSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxDeleteSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxAddImageToSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxAddTableToSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxAddShapeToSlideTool(thread_id=thread_id, emit_events=emit_events),
        SandboxGetPresentationInfoTool(thread_id=thread_id, emit_events=emit_events),
    ]
